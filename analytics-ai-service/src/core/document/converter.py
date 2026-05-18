"""
Office document → Markdown converters.

Both converters write the resulting Markdown to disk and return the output
path, so the caller can feed the file straight into PageIndex's vendored
markdown engine (md_to_tree). The two libraries used here are pure-Python
(no LibreOffice / Java required), keeping the runtime image small.

Conversion choices:
  - .docx: mammoth preserves Word's Heading 1/2/3 paragraph styles as
           Markdown # / ## / ### so PageIndex can build a tree that
           matches the document's actual outline. Other formatting
           (bold, italic, lists, links, tables) is preserved as standard
           Markdown.
  - .pptx: python-pptx exposes raw shapes; we synthesize Markdown so each
           slide becomes a `# Slide N: <title>` heading with bullet text
           and (optional) `## Notes` subsection. This gives PageIndex a
           natural per-slide hierarchy.

Both formats fail loudly with ConversionError if the source can't be
parsed (corrupt file, password-protected, etc.) so the caller can mark
the document as failed with a useful message.
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger("analytics-service")


class ConversionError(Exception):
    """Raised when a document can't be converted to Markdown."""


def docx_to_markdown(src: Path, dst: Path) -> Path:
    """Convert a .docx file to Markdown, writing the result to `dst`.

    Returns the path to the written .md file. mammoth `messages` (warnings
    about unhandled styles, missing images, etc.) are logged but never
    raised — they are informational, not fatal.
    """
    import mammoth

    try:
        with open(src, "rb") as f:
            result = mammoth.convert_to_markdown(f)
    except Exception as e:
        raise ConversionError(
            f"Failed to read .docx file '{src.name}': {e!r}"
        ) from e

    if result.messages:
        warnings = [m.message for m in result.messages if m.type == "warning"]
        if warnings:
            logger.info(
                f"mammoth produced {len(warnings)} warning(s) converting "
                f"{src.name}: {warnings[:3]}{'...' if len(warnings) > 3 else ''}"
            )

    markdown = (result.value or "").strip()
    if not markdown:
        raise ConversionError(
            f"'{src.name}' contains no extractable text"
        )

    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(markdown, encoding="utf-8")
    return dst


def pptx_to_markdown(src: Path, dst: Path) -> Path:
    """Convert a .pptx file to Markdown, writing the result to `dst`.

    Each slide becomes one `# Slide N: <title>` heading. Body text from
    every text frame on the slide is emitted as a bullet list. Speaker
    notes (if any) go into a `## Notes` subsection so they're still
    indexed but separable from the slide content itself.

    Quality improvements vs naive extraction:
      - Title fallback: when the title shape is missing/empty (very common
        in modern decks), we fall back to the first non-trivial text shape
        on the slide so the heading is informative ("Slide 3: Pricing")
        instead of a generic placeholder ("Slide 3: Slide 3").
      - Noise filtering: we drop tokens that are obviously template
        artefacts — single characters, pure symbols/arrows, and short
        strings that repeat across most slides (logo names, page chrome).
    """
    from pptx import Presentation

    try:
        prs = Presentation(str(src))
    except Exception as e:
        raise ConversionError(
            f"Failed to read .pptx file '{src.name}': {e!r}"
        ) from e

    # First pass — collect every non-title text fragment so we can spot
    # strings that repeat across most slides (those are almost always
    # boilerplate: logos, page footers, brand handles).
    per_slide_fragments: list[list[str]] = []
    for slide in prs.slides:
        per_slide_fragments.append(_collect_slide_fragments(slide))
    boilerplate = _detect_boilerplate(per_slide_fragments)

    lines: list[str] = []
    for idx, (slide, fragments) in enumerate(
        zip(prs.slides, per_slide_fragments), start=1
    ):
        # Clean fragments first; we'll use them both to pick a title
        # fallback and to emit the bullet list.
        clean = [f for f in fragments if _keep_fragment(f, boilerplate)]

        title_text = _extract_slide_title(slide)
        title_from_fallback = False
        if not title_text and clean:
            title_text = clean[0]
            title_from_fallback = True
        if not title_text:
            title_text = f"Slide {idx}"

        # Use idx in the heading so we keep a stable ordering even when
        # two slides happen to share a title (common in deck templates).
        lines.append(f"# Slide {idx}: {title_text}")
        lines.append("")

        # Avoid emitting the fallback title again as the first bullet.
        body = clean[1:] if title_from_fallback else clean
        for text in body:
            lines.append(f"- {text}")
        lines.append("")

        # Speaker notes as a sub-section so retrieval can find them but
        # they don't pollute the slide's own content.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("## Notes")
                lines.append("")
                lines.append(notes)
                lines.append("")

    markdown = "\n".join(lines).strip()
    if not markdown:
        raise ConversionError(
            f"'{src.name}' contains no extractable text "
            f"(deck may be image-only)"
        )

    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_text(markdown, encoding="utf-8")
    return dst


def _collect_slide_fragments(slide) -> list[str]:
    """Return every non-title paragraph text from a slide, in reading order."""
    title_shape = getattr(slide.shapes, "title", None)
    fragments: list[str] = []
    for shape in slide.shapes:
        if shape is title_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                fragments.append(text)
    return fragments


# Single-character symbols / arrows / bullets we never want as their own bullet.
_NOISE_SINGLE_CHARS = set("→←↑↓⇒⇐⋯…·•●○▪▫■□◆◇★☆⭐")


def _keep_fragment(text: str, boilerplate: set[str]) -> bool:
    """Filter rules: drop obvious template noise but keep real content."""
    if text in boilerplate:
        return False
    # Single-character text is almost always a glyph placeholder (logo
    # initial, arrow icon, decorative bullet). Numbers are an exception
    # because page numbers / pricing tiers can be a single digit.
    if len(text) == 1 and not text.isdigit():
        return False
    # Pure symbols (arrows, icons) carry no semantic content.
    if all(ch in _NOISE_SINGLE_CHARS or not ch.isalnum() for ch in text):
        return False
    return True


def _detect_boilerplate(per_slide: list[list[str]]) -> set[str]:
    """Find short strings that repeat across the majority of slides.

    Brand names, page chrome, and footer handles tend to be present on
    almost every slide and add zero retrieval value. Anything appearing
    on > 50% of slides AND under 40 characters qualifies. We keep the
    threshold tight on length so we never drop genuine repeated content
    like a heading that recurs in a recap deck.
    """
    if len(per_slide) < 3:
        # Not enough slides to reliably distinguish boilerplate from
        # legitimate repeats; bail out.
        return set()

    counts: dict[str, int] = {}
    for fragments in per_slide:
        for f in set(fragments):  # dedupe within a slide first
            counts[f] = counts.get(f, 0) + 1

    threshold = len(per_slide) / 2
    return {
        text
        for text, n in counts.items()
        if n > threshold and len(text) < 40
    }


def _extract_slide_title(slide) -> str | None:
    """Best-effort: return the slide title text, or None if unavailable."""
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is None:
        return None
    if not getattr(title_shape, "has_text_frame", False):
        return None
    text = title_shape.text_frame.text.strip()
    return text or None
